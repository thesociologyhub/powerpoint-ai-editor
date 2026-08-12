from flask import Flask, request, send_file, jsonify
from pptx import Presentation
from openai import OpenAI
import io
import os
import json

app = Flask(__name__)

client = OpenAI(
    api_key=os.environ.get("OPENAI_API_KEY"),
    timeout=90.0,
    max_retries=2,
)

SYSTEM_PROMPT = """
You are The Sociology Hub PowerPoint Proofreading Assistant.

Make ONLY safe editorial corrections.

You MAY correct:
- spelling
- grammar
- punctuation
- obvious typos
- duplicated words
- capitalisation consistency

Use UK English.

You MUST NOT change:
- sociology content
- sociological terminology
- curriculum content
- teaching activities
- exam questions
- learning objectives
- pedagogical decisions
- factual claims
- examples
- meaning or tone

Preserve wording as closely as possible.

You will receive a numbered JSON list of PowerPoint paragraphs.

Return ONLY valid JSON in exactly this format:

{"items":[{"id":0,"text":"corrected text"}]}

Return every supplied id exactly once and in the same order.
"""


@app.route("/")
def home():
    return "PowerPoint AI Editor is running!"


def collect_paragraphs(presentation):
    paragraphs = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue

            for paragraph in shape.text_frame.paragraphs:
                text = "".join(run.text for run in paragraph.runs)

                if text.strip():
                    paragraphs.append({
                        "id": len(paragraphs),
                        "text": text,
                        "paragraph": paragraph
                    })

    return paragraphs


def proofread_batch(items):
    payload = [
        {"id": item["id"], "text": item["text"]}
        for item in items
    ]

    response = client.responses.create(
        model="gpt-5-mini",
        instructions=SYSTEM_PROMPT,
        input=json.dumps(payload, ensure_ascii=False),
    )

    result = json.loads(response.output_text)

    return result["items"]


@app.route("/edit", methods=["POST"])
def edit_powerpoint():
    try:
        if "file" not in request.files:
            return jsonify({"error": "No PowerPoint file received"}), 400

        if not os.environ.get("OPENAI_API_KEY"):
            return jsonify({"error": "OPENAI_API_KEY is not configured"}), 500

        uploaded_file = request.files["file"]
        presentation = Presentation(uploaded_file)

        paragraphs = collect_paragraphs(presentation)

        batch_size = 40

        for start in range(0, len(paragraphs), batch_size):
            batch = paragraphs[start:start + batch_size]
            corrections = proofread_batch(batch)

            correction_map = {
                item["id"]: item["text"]
                for item in corrections
            }

            for item in batch:
                corrected = correction_map.get(item["id"], item["text"])
                paragraph = item["paragraph"]

                if corrected != item["text"] and paragraph.runs:
                    paragraph.runs[0].text = corrected

                    for run in paragraph.runs[1:]:
                        run.text = ""

        output = io.BytesIO()
        presentation.save(output)
        output.seek(0)

        original_name = uploaded_file.filename or "presentation.pptx"

        if original_name.lower().endswith(".pptx"):
            new_name = original_name[:-5] + " - Corrected.pptx"
        else:
            new_name = "Corrected PowerPoint.pptx"

        return send_file(
            output,
            as_attachment=True,
            download_name=new_name,
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

    except Exception as e:
        print(f"ERROR: {repr(e)}", flush=True)
        return jsonify({"error": str(e)}), 500
